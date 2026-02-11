from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_premium_presentation():
    prs = Presentation()
    
    # Define Colors
    COLOR_SAFFRON = RGBColor(255, 153, 51)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_NAVY = RGBColor(0, 0, 128)
    COLOR_DARK_BG = RGBColor(10, 10, 25)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_top_bar(slide):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_SAFFRON
        shape.line.visible = False

    # Slide 1: Title Slide (The Reveal)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_DARK_BG)
    add_top_bar(slide)
    
    # Big 'Y' Logo Background (Circle)
    logo_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, (prs.slide_width - Inches(2.5))/2, Inches(1.2), Inches(2.5), Inches(2.5))
    logo_bg.fill.solid()
    logo_bg.fill.fore_color.rgb = COLOR_SAFFRON
    logo_bg.line.visible = False
    
    logo_text = slide.shapes.add_textbox((prs.slide_width - Inches(2.5))/2, Inches(1.5), Inches(2.5), Inches(2))
    tf = logo_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Y"
    p.font.size = Pt(100)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(0, Inches(4), prs.slide_width, Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PLATFORM Y"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLOR_SAFFRON
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(0, Inches(5.2), prs.slide_width, Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "RECLAIMING INDIA'S DIGITAL SOVEREIGNTY"
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: The Crisis (Gradients & Cards)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_DARK_BG)
    add_top_bar(slide)
    
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
    p = header.text_frame.paragraphs[0]
    p.text = "THE CRISIS: INFORMATION WARFARE"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_SAFFRON

    points = [
        {"title": "Foreign Control", "text": "National security decisions held hostage by global tech giants."},
        {"title": "Anti-India Bias", "text": "Unfiltered negative narratives portrayal of 1.4B people."},
        {"title": "Data Hemorrhage", "text": "Critical data stored in offshore servers beyond Indian law."}
    ]

    for i, pt in enumerate(points):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + i*3.1), Inches(2), Inches(2.9), Inches(3.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(15, 15, 45)
        shape.line.color.rgb = COLOR_SAFFRON
        shape.line.width = Pt(2.5)
        
        tb = slide.shapes.add_textbox(Inches(0.6 + i*3.1), Inches(2.3), Inches(2.7), Inches(3))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = pt["title"]
        p1.font.bold = True
        p1.font.size = Pt(20)
        p1.font.color.rgb = COLOR_SAFFRON
        
        p2 = tf.add_paragraph()
        p2.text = "\n" + pt["text"]
        p2.font.size = Pt(17)
        p2.font.color.rgb = COLOR_WHITE

    # Slide 3: Sovereign Moderation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_DARK_BG)
    add_top_bar(slide)

    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
    p = header.text_frame.paragraphs[0]
    p.text = "SOVEREIGN CONTROL: THE GOI PORTAL"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_SAFFRON

    # Bullet points with saffron checks
    features = [
        "Sovereign 'Super Access' for Govt Agencies",
        "Global Content Deletion Power (Sovereign Right)",
        "Zero-Permission Neutralization of Fake News",
        "Independent Indian Infrastructure"
    ]
    
    for i, f in enumerate(features):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2 + i*1.1), prs.slide_width - Inches(1), Inches(0.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(25, 25, 60)
        box.line.visible = False
        
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "  ▶  " + f
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.LEFT

    # Slide 4: World Class Technology
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_DARK_BG)
    add_top_bar(slide)
    
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
    p = header.text_frame.paragraphs[0]
    p.text = "THE TECH STACK: PLATFORM Y"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_SAFFRON

    stacks = [
        {"name": "GO LANG", "desc": "Power for 2B Citizens"},
        {"name": "FLUTTER", "desc": "Ultra-Smooth UI"},
        {"name": "NIC CLOUD", "desc": "Sovereign Hosting"}
    ]

    for i, s in enumerate(stacks):
        rect = slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_EXTRACT, Inches(0.5 + i*3.2), Inches(2.5), Inches(2.5), Inches(2.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLOR_SAFFRON
        rect.line.visible = False
        
        tb = slide.shapes.add_textbox(Inches(0.5 + i*3.2), Inches(3.2), Inches(2.5), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = s["name"]
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_NAVY
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tb.text_frame.add_paragraph()
        p2.text = s["desc"]
        p2.font.size = Pt(16)
        p2.font.color.rgb = COLOR_DARK_BG
        p2.alignment = PP_ALIGN.CENTER

    # Save
    import os
    save_path = "Documentation/09_Presentations/Platform_Y_World_Class.pptx"
    os.makedirs(os.path.dirname(save_path), exist_ok=True)
    prs.save(save_path)
    print(f"Presentation saved to {save_path}")

if __name__ == "__main__":
    create_premium_presentation()
